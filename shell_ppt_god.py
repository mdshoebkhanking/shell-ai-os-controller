import os
import json
import logging
import asyncio
from datetime import datetime
from typing import List, Optional
from shell_safe_executor import god_tier_tool as function_tool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Shell AI infrastructure
from shell_config import config
from shell_logger import get_logger

logger = get_logger("shell_ppt_god")

# --- GROQ SETUP ---
try:
    from groq import Groq
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False

client = None
if GROQ_AVAILABLE and config.get_str("GROQ_API_KEY"):
    try:
        client = Groq(api_key=config.get_str("GROQ_API_KEY"))
    except Exception as e:
        logger.error(f"Groq Client Init Failed: {e}")

# --- STYLE COLOR SCHEMES ---
STYLE_SCHEMES = {
    "professional": {
        "title_color": RGBColor(0x1A, 0x1A, 0x2E),      # Dark navy
        "subtitle_color": RGBColor(0x4A, 0x4A, 0x6A),    # Muted purple-gray
        "heading_color": RGBColor(0x2C, 0x3E, 0x50),     # Dark blue-gray
        "body_color": RGBColor(0x33, 0x33, 0x33),        # Dark gray
        # Background hint: White/Light gray - classic corporate look
    },
    "modern": {
        "title_color": RGBColor(0x6C, 0x5C, 0xE7),      # Purple
        "subtitle_color": RGBColor(0xA2, 0x9B, 0xFE),    # Light purple
        "heading_color": RGBColor(0x00, 0xB8, 0x94),     # Teal
        "body_color": RGBColor(0x2D, 0x3A, 0x4A),        # Dark slate
        # Background hint: White with accent colors - trendy startup vibe
    },
    "minimal": {
        "title_color": RGBColor(0x22, 0x22, 0x22),      # Near black
        "subtitle_color": RGBColor(0x88, 0x88, 0x88),    # Medium gray
        "heading_color": RGBColor(0x33, 0x33, 0x33),     # Dark gray
        "body_color": RGBColor(0x55, 0x55, 0x55),        # Gray
        # Background hint: Pure white - clean and distraction-free
    },
    "dark": {
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),      # White
        "subtitle_color": RGBColor(0xBB, 0xBB, 0xBB),    # Light gray
        "heading_color": RGBColor(0x00, 0xD2, 0xFF),     # Cyan accent
        "body_color": RGBColor(0xDD, 0xDD, 0xDD),        # Off-white
        # Background hint: Dark/black - use dark slide backgrounds for best effect
    },
}

async def generate_ppt_content(topic: str, num_slides: int) -> dict:
    """Uses Groq to generate structured JSON content for the PPT."""
    if not client:
        logger.error("Groq client not initialized.")
        return None

    prompt = f"""
    Create a professional PowerPoint presentation outline on the topic: '{topic}'.
    It should have exactly {num_slides} slides.

    Return a VALID JSON object strictly in this format:
    {{
      "title": "Presentation Title",
      "subtitle": "Subtitle or Slogan",
      "slides": [
        {{
          "layout": "bullet",
          "title": "Slide Title",
          "content": ["Point 1", "Point 2", "Point 3"]
        }}
      ]
    }}

    Layout types: 'title', 'bullet', 'section'.
    First slide should be the main title slide (handled by title/subtitle fields), so 'slides' list is for subsequent slides.
    Do not use markdown formatting. Return clean JSON only.
    """

    try:
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": prompt,
                }
            ],
            model="llama-3.3-70b-versatile",
            response_format={"type": "json_object"}
        )
        text = response.choices[0].message.content.strip()

        # Cleanup if markdown blocks exist despite instruction
        if text.startswith("```json\n"): text = text[8:]
        elif text.startswith("```json"): text = text[7:]
        elif text.startswith("```"): text = text[3:]
        if text.endswith("```"): text = text[:-3]

        return json.loads(text.strip())
    except Exception as e:
        logger.error(f"Content Gen with Groq failed: {e}")
        return None

def _apply_style_to_text(text_frame, color: RGBColor, font_size: int = None):
    """Helper: Apply color and optional font size to all paragraphs in a text frame."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color
            if font_size:
                run.font.size = Pt(font_size)

@function_tool
async def create_presentation_tool(topic: str, num_slides: int = 5, filename: str = None, style: str = "professional") -> str:
    """
    Creates a PowerPoint presentation on the given topic using AI.

    Args:
        topic (str): The subject of the presentation.
        num_slides (int): Number of content slides (default 5).
        filename (str): Optional filename.
        style (str): Visual style - professional, modern, minimal, dark (default professional).
    """
    try:
        # Validate style
        style = style.lower() if style else "professional"
        if style not in STYLE_SCHEMES:
            style = "professional"
        scheme = STYLE_SCHEMES[style]

        logger.info(f"🎨 Generating PPT on '{topic}' ({num_slides} slides, style: {style})...")

        # 1. Generate Content
        content_data = await generate_ppt_content(topic, num_slides)

        if not content_data:
            return "❌ Presentation content generate nahi ho paya. AI Brain offline ya error hai."

        # 2. Create Presentation
        prs = Presentation()

        # --- Slide 1: Title Slide ---
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = content_data.get("title", topic)
        subtitle.text = content_data.get("subtitle", "Generated by Shell AI")

        # Apply style colors to title slide
        _apply_style_to_text(title.text_frame, scheme["title_color"])
        _apply_style_to_text(subtitle.text_frame, scheme["subtitle_color"])

        # --- Content Slides ---
        for i, slide_data in enumerate(content_data.get("slides", [])):
            layout_type = slide_data.get("layout", "bullet")

            # Choose Layout (1 = Title + Content)
            slide_layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(slide_layout)

            # Set Title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get("title", f"Slide {i+1}")
                _apply_style_to_text(slide.shapes.title.text_frame, scheme["heading_color"])

            # Set Content (Bullets)
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame

                points = slide_data.get("content", [])
                if points:
                    tf.text = points[0] # First point
                    for point in points[1:]:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0

                # Apply body text color
                _apply_style_to_text(tf, scheme["body_color"])

        # 3. Save
        workspace_dir = os.path.join(os.getcwd(), "shell_projects", "presentations")
        if not os.path.exists(workspace_dir):
            os.makedirs(workspace_dir)

        if not filename:
            clean_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '_')).replace(' ', '_')
            filename = f"{clean_topic}.pptx"

        if not filename.endswith(".pptx"):
            filename += ".pptx"

        save_path = os.path.join(workspace_dir, filename)
        prs.save(save_path)

        logger.info(f"💾 PPT Saved: {save_path}")

        # Get slide count and file size
        total_slides = len(prs.slides)
        file_size_bytes = os.path.getsize(save_path)
        if file_size_bytes >= 1024 * 1024:
            file_size_str = f"{file_size_bytes / (1024 * 1024):.2f} MB"
        else:
            file_size_str = f"{file_size_bytes / 1024:.2f} KB"

        # 4. Open it
        os.startfile(save_path)

        return (
            f"✅ Presentation ban gayi Sir! Style: {style.capitalize()}\n"
            f"📊 Total Slides: {total_slides} | File Size: {file_size_str}\n"
            f"📁 Saved to: {filename}\n"
            f"📍 Path: {save_path}"
        )

    except Exception as e:
        logger.error(f"PPT Creation Error: {e}")
        return f"❌ PPT banane mein error aaya: {str(e)}"


@function_tool
async def list_presentations_tool() -> str:
    """
    Lists all .pptx presentation files in the shell_projects/presentations/ directory.
    Shows name, size, modified date, and slide count for each file.
    """
    try:
        workspace_dir = os.path.join(os.getcwd(), "shell_projects", "presentations")

        if not os.path.exists(workspace_dir):
            return "📂 Presentations folder abhi exist nahi karta Sir. Pehle ek presentation banayein!"

        pptx_files = [f for f in os.listdir(workspace_dir) if f.lower().endswith(".pptx")]

        if not pptx_files:
            return "📂 Koi presentation file nahi mili Sir. Pehle ek banayein `create_presentation_tool` se!"

        result_lines = [f"📂 --- Presentations List ({len(pptx_files)} files) ---\n"]

        for fname in sorted(pptx_files):
            fpath = os.path.join(workspace_dir, fname)
            try:
                # File size
                size_bytes = os.path.getsize(fpath)
                if size_bytes >= 1024 * 1024:
                    size_str = f"{size_bytes / (1024 * 1024):.2f} MB"
                else:
                    size_str = f"{size_bytes / 1024:.2f} KB"

                # Modified date
                mod_time = os.path.getmtime(fpath)
                mod_date = datetime.fromtimestamp(mod_time).strftime("%Y-%m-%d %H:%M")

                # Slide count
                try:
                    prs = Presentation(fpath)
                    slide_count = len(prs.slides)
                except Exception:
                    slide_count = "?"

                result_lines.append(
                    f"📄 {fname}\n"
                    f"   Size: {size_str} | Slides: {slide_count} | Modified: {mod_date}"
                )
            except Exception as e:
                result_lines.append(f"📄 {fname} — Error reading details: {e}")

        return "\n".join(result_lines)

    except Exception as e:
        logger.error(f"List Presentations Error: {e}")
        return f"❌ Presentations list karne mein error: {str(e)}"


@function_tool
async def add_slide_to_ppt_tool(filepath: str, slide_title: str, bullet_points: str) -> str:
    """
    Opens an existing .pptx file and adds a new slide with title and bullet points.

    Args:
        filepath (str): Path to the existing .pptx file.
        slide_title (str): Title for the new slide.
        bullet_points (str): Comma-separated bullet points for the slide content.
    """
    try:
        # Resolve path - check absolute or relative to presentations dir
        if not os.path.isabs(filepath):
            workspace_dir = os.path.join(os.getcwd(), "shell_projects", "presentations")
            filepath = os.path.join(workspace_dir, filepath)

        if not os.path.exists(filepath):
            return f"❌ File nahi mili Sir: {filepath}\nPehle check karein path sahi hai ya nahi."

        if not filepath.lower().endswith(".pptx"):
            return "❌ Ye .pptx file nahi hai Sir. Sirf PowerPoint files supported hain."

        # Open existing presentation
        prs = Presentation(filepath)

        # Add new slide with Title + Content layout
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_title

        # Set bullet points
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame

            points = [p.strip() for p in bullet_points.split(",") if p.strip()]
            if points:
                tf.text = points[0]
                for point in points[1:]:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0

        # Save
        prs.save(filepath)

        total_slides = len(prs.slides)
        fname = os.path.basename(filepath)

        return (
            f"✅ Naya slide add ho gaya Sir!\n"
            f"📄 File: {fname}\n"
            f"📊 Total Slides ab: {total_slides}\n"
            f"📝 Title: {slide_title}"
        )

    except Exception as e:
        logger.error(f"Add Slide Error: {e}")
        return f"❌ Slide add karne mein error: {str(e)}"
